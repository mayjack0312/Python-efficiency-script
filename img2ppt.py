import os
import pptx
from pptx.util import Inches

ppt_filename = input('输入目标ppt文件名(无需后缀)：')
full_ppt_filename = '{}.{}'.format(ppt_filename,'pptx')
ppt_file = pptx.Presentation()

pic_files = [fn for fn in os.listdir() if fn.endswith('.png')]

# 按图片编号顺序导入
for fn in sorted(pic_files, key=lambda item:int(item[:item.rindex('.')])):
    slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])

    # 为PPTX文件当前幻灯片中第一个文本框设置文字，本文代码中可忽略
    # slide.shapes.placeholders[0].text = fn[:fn.rindex('.')]

    # 导入并为当前幻灯片添加图片，起始位置和尺寸可修改
    slide.shapes.add_picture(fn, Inches(0), Inches(0), Inches(10), Inches(7.5))
    
ppt_file.save(full_ppt_filename)
