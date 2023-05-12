from pptx import Presentation
import os
import sys

ppt_root = jpg_root = sys.path[0]
print(ppt_root)
i = 0

def ppt2png(pptFileName):
    # 实例化ppt对象
    prs = Presentation(pptFileName)  # PPT文件路径
    # 遍历

    for slide in prs.slides:
        for obj in slide.shapes:
            # try/except 因为非图片元素没有Image方法会弹出异常.
            try:
                global i
                # 获取二进制字符流
                imdata = obj.image.blob
                # 判断文件后缀类型
                imagetype = obj.image.content_type
                typekey = imagetype.find('/') + 1
                i += 1
                imtype = imagetype[typekey:]

                # 创建image文件夹保存抽出图片
                path = jpg_root + "/image/"

                if not os.path.exists(path):
                    os.makedirs(path)
                # 图片生成
                obj.name = "Picture" + str(i)
                image_file = path + obj.name + "." + imtype

                file_str = open(image_file, 'wb')
                file_str.write(imdata)
                file_str.close()
            except:
                pass

for fn in (fns for fns in os.listdir(ppt_root) if fns.endswith(('.ppt', 'pptx'))):
    ppt2png(fn)